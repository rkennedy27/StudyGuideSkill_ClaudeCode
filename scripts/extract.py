"""
Hybrid extraction for lecture sources (PDF, PPTX).

Strategy:
  1. Pull text from every page/slide via PyMuPDF (PDF) or python-pptx (PPTX) — cheap.
  2. For each page, decide if it's "vision-needed" using two heuristics:
       - extracted text < MIN_WORDS (likely image-only slide), OR
       - page contains embedded images that occupy a meaningful area
  3. Render only the flagged pages as PNG to assets/<source>/slide_NNN.png.
  4. Emit one <source>.txt of full text and one <source>.flags.json listing
     flagged page numbers + their PNG paths.

The downstream extraction agent reads the .txt for everything it can,
and Read()s only the flagged PNGs — keeping image token spend minimal.

Usage:
    python extract.py <slides_dir> <output_dir>

Requires: pymupdf, python-pptx (only if .pptx files present)
"""

import json
import os
import sys
from pathlib import Path

MIN_WORDS = 50
IMAGE_AREA_RATIO = 0.30  # flag page if embedded images cover >30% of page area
RENDER_DPI = 180


def extract_pdf(pdf_path: Path, out_text: Path, assets_dir: Path) -> dict:
    import fitz

    doc = fitz.open(pdf_path)
    text_parts = []
    flags = []

    for i, page in enumerate(doc):
        page_num = i + 1
        text = page.get_text() or ""
        word_count = len(text.split())

        page_area = page.rect.width * page.rect.height
        image_area = 0.0
        try:
            for img in page.get_images(full=True):
                xref = img[0]
                for rect in page.get_image_rects(xref):
                    image_area += rect.width * rect.height
        except Exception:
            pass
        image_ratio = image_area / page_area if page_area else 0.0

        sparse = word_count < MIN_WORDS
        image_heavy = image_ratio > IMAGE_AREA_RATIO

        text_parts.append(f"\n\n===== PAGE {page_num} =====\n{text}".rstrip())

        if sparse or image_heavy:
            png_name = f"slide_{page_num:03d}.png"
            png_path = assets_dir / png_name
            assets_dir.mkdir(parents=True, exist_ok=True)
            page.get_pixmap(dpi=RENDER_DPI).save(str(png_path))
            flags.append({
                "page": page_num,
                "reason": "sparse_text" if sparse else "image_heavy",
                "word_count": word_count,
                "image_ratio": round(image_ratio, 3),
                "image": str(png_path.relative_to(assets_dir.parent.parent)),
            })

    doc.close()
    out_text.write_text("".join(text_parts), encoding="utf-8")
    return {"pages": len(text_parts), "flagged": flags}


def extract_pptx(pptx_path: Path, out_text: Path, assets_dir: Path) -> dict:
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        print(f"  python-pptx not installed; skipping {pptx_path.name}", file=sys.stderr)
        return {"pages": 0, "flagged": [], "skipped": True}

    prs = Presentation(pptx_path)
    text_parts = []
    flags = []

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    slide_area = slide_w * slide_h if slide_w and slide_h else 1

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        chunks = []
        image_area = 0

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        chunks.append(line)
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image_area += (shape.width or 0) * (shape.height or 0)
                except Exception:
                    pass

        text = "\n".join(chunks)
        word_count = len(text.split())
        image_ratio = image_area / slide_area
        sparse = word_count < MIN_WORDS
        image_heavy = image_ratio > IMAGE_AREA_RATIO

        text_parts.append(f"\n\n===== SLIDE {slide_num} =====\n{text}".rstrip())

        if sparse or image_heavy:
            # PPTX has no built-in render; rely on user converting to PDF first,
            # OR fall back to skipping render and just flagging.
            flags.append({
                "page": slide_num,
                "reason": "sparse_text" if sparse else "image_heavy",
                "word_count": word_count,
                "image_ratio": round(image_ratio, 3),
                "image": None,  # no render available for pptx without LibreOffice
                "note": "Convert PPTX to PDF for image rendering, or read source PPTX directly.",
            })

    out_text.write_text("".join(text_parts), encoding="utf-8")
    return {"pages": len(text_parts), "flagged": flags}


def main():
    if len(sys.argv) < 3:
        print("Usage: python extract.py <slides_dir> <output_dir>", file=sys.stderr)
        sys.exit(2)

    slides_dir = Path(sys.argv[1]).resolve()
    out_dir = Path(sys.argv[2]).resolve()
    sources_text_dir = out_dir / "sources-text"
    assets_root = out_dir / "assets"
    sources_text_dir.mkdir(parents=True, exist_ok=True)

    files = sorted(
        [p for p in slides_dir.iterdir() if p.suffix.lower() in (".pdf", ".pptx")]
    )
    if not files:
        print(f"No PDF/PPTX files in {slides_dir}", file=sys.stderr)
        sys.exit(1)

    summary = {}
    for src in files:
        base = src.stem.replace(" - Tagged", "").strip()
        out_text = sources_text_dir / f"{base}.txt"
        out_flags = sources_text_dir / f"{base}.flags.json"
        assets_dir = assets_root / base

        print(f"Extracting {src.name} ...")
        if src.suffix.lower() == ".pdf":
            result = extract_pdf(src, out_text, assets_dir)
        else:
            result = extract_pptx(src, out_text, assets_dir)

        out_flags.write_text(json.dumps(result, indent=2), encoding="utf-8")
        summary[base] = {
            "pages": result["pages"],
            "flagged_count": len(result["flagged"]),
        }
        print(f"  -> {result['pages']} pages, {len(result['flagged'])} flagged for vision")

    (out_dir / "extraction-summary.json").write_text(
        json.dumps(summary, indent=2), encoding="utf-8"
    )
    print("\nDone.")


if __name__ == "__main__":
    main()
